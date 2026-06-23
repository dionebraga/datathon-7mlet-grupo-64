"""
Atualiza o PPTX da Adaptive Offers Platform com a nova paleta de cores
e adiciona slides para MLflow e API Docs.

Nova paleta:
  BG/fundo    #000000  preto
  PANEL       #030D24  navy escuro
  VIOLET      #0033CC  azul royal
  CYAN        #1A6FFF  azul elétrico
  GREEN       #1A9E1A  verde vivo
  GOLD        #FFC200  amarelo ouro
  AMBER       #FF9A00  laranja quente
  RED         #E84000  laranja-vermelho
  TEXT        #EDEDED  branco suave
  MUTED       #8899BB  cinza-azul
"""

from __future__ import annotations

import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.dml.color import RGBColor
from pptx.util import Pt
import copy

# ── Palette ──────────────────────────────────────────────────────────────────
BG     = RGBColor(0x00, 0x00, 0x00)
PANEL  = RGBColor(0x03, 0x0D, 0x24)
VIOLET = RGBColor(0x00, 0x33, 0xCC)
CYAN   = RGBColor(0x1A, 0x6F, 0xFF)
GREEN  = RGBColor(0x1A, 0x9E, 0x1A)
GOLD   = RGBColor(0xFF, 0xC2, 0x00)
AMBER  = RGBColor(0xFF, 0x9A, 0x00)
RED    = RGBColor(0xE8, 0x40, 0x00)
TEXT   = RGBColor(0xED, 0xED, 0xED)
MUTED  = RGBColor(0x88, 0x99, 0xBB)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LIME   = RGBColor(0xA0, 0xC8, 0x30)

# Map old accent colors to new palette (approximate substitutions)
OLD_TO_NEW: dict[str, RGBColor] = {
    # old purples / accents → royal blue
    "5B4EDB": VIOLET, "6366F1": VIOLET, "7C3AED": VIOLET,
    "4F46E5": VIOLET, "8B5CF6": CYAN,
    # old teals / cyans
    "06B6D4": CYAN, "0EA5E9": CYAN, "38BDF8": CYAN,
    "00B4D8": CYAN, "0077B6": CYAN,
    # old greens
    "10B981": GREEN, "34D399": GREEN, "22C55E": GREEN,
    "86EFAC": GREEN, "4ADE80": GREEN,
    # old yellow/amber
    "F59E0B": GOLD, "FBBF24": GOLD, "FCD34D": GOLD,
    "EAB308": GOLD, "CA8A04": AMBER,
    # old reds/orange
    "EF4444": RED, "F97316": RED, "DC2626": RED,
    "E11D48": RED,
    # old dark bg → new panel
    "0A0A0A": PANEL, "111111": PANEL, "0D0D0D": PANEL,
    "1A1A2E": PANEL, "16213E": PANEL, "0F3460": PANEL,
    # old grey → muted
    "94A3B8": MUTED, "64748B": MUTED,
    # white-ish text stays white/near-white
    "F8FAFC": TEXT, "F1F5F9": TEXT,
}


def hex_to_rgb(hex_str: str) -> RGBColor | None:
    h = hex_str.upper().strip("#")
    if len(h) == 6:
        try:
            return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
        except ValueError:
            return None
    return None


def rgb_to_hex(color: RGBColor) -> str:
    return f"{color[0]:02X}{color[1]:02X}{color[2]:02X}"


def remap_color(rgb: RGBColor) -> RGBColor | None:
    """If the color is in the mapping table, return the new color."""
    key = rgb_to_hex(rgb)
    return OLD_TO_NEW.get(key, None)


def update_color_element(elem):
    """Walk XML element and update srgbClr values."""
    from lxml import etree
    NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
    for srgb in elem.iter(f"{{{NS}}}srgbClr"):
        val = srgb.get("val", "")
        mapped = OLD_TO_NEW.get(val.upper(), None)
        if mapped:
            srgb.set("val", rgb_to_hex(mapped).lower())


def process_shape(shape):
    """Update colors in shape fills, outlines, and text."""
    try:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.color.type is not None:
                        try:
                            c = run.font.color.rgb
                            mapped = remap_color(c)
                            if mapped:
                                run.font.color.rgb = mapped
                        except Exception:
                            pass
    except Exception:
        pass
    # update XML-level color refs
    try:
        update_color_element(shape._element)
    except Exception:
        pass


def update_slide_background(slide, color: RGBColor = PANEL):
    """Set solid background fill on a slide."""
    from lxml import etree
    from pptx.oxml.ns import qn
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def make_title_bar(slide, prs, title_text: str, subtitle: str = ""):
    """Add branded title strip at top of slide."""
    W = prs.slide_width
    # top accent bar (thin gold line)
    bar = slide.shapes.add_shape(1, 0, Inches(0.0), W, Inches(0.07))  # MSO_SHAPE_TYPE.RECTANGLE=1
    bar.fill.solid()
    bar.fill.fore_color.rgb = GOLD
    bar.line.fill.background()

    # title box
    txb = slide.shapes.add_textbox(Inches(0.45), Inches(0.12), W - Inches(0.9), Inches(0.65))
    tf = txb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.bold = True
    p.font.size = Pt(24)
    p.font.color.rgb = TEXT
    p.font.name = "Inter"

    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.45), Inches(0.75), W - Inches(0.9), Inches(0.35))
        stf = sub.text_frame
        sp = stf.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(13)
        sp.font.color.rgb = MUTED
        sp.font.name = "Inter"


def add_bullet_card(slide, left, top, width, height, title, bullets, accent=CYAN):
    """Add a card-style box with title and bullet list."""
    # card bg
    card = slide.shapes.add_shape(1, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0x05, 0x14, 0x2C)  # slightly lighter than PANEL
    card.line.color.rgb = accent
    card.line.width = Pt(1.2)

    # title
    txb = slide.shapes.add_textbox(left + Inches(0.18), top + Inches(0.12), width - Inches(0.36), Inches(0.35))
    tf = txb.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = Pt(13)
    p.font.color.rgb = accent
    p.font.name = "Inter"

    # bullet list
    bxb = slide.shapes.add_textbox(left + Inches(0.22), top + Inches(0.5), width - Inches(0.44), height - Inches(0.65))
    btf = bxb.text_frame
    btf.word_wrap = True
    for i, bullet in enumerate(bullets):
        bp = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
        bp.text = f"▸  {bullet}"
        bp.font.size = Pt(11)
        bp.font.color.rgb = TEXT
        bp.font.name = "Inter"
        bp.space_after = Pt(4)


def add_mlflow_slide(prs: Presentation):
    """Add a slide explaining MLflow experiment tracking integration."""
    slide_layout = prs.slide_layouts[0]  # only layout available
    slide = prs.slides.add_slide(slide_layout)
    W, H = prs.slide_width, prs.slide_height

    update_slide_background(slide, PANEL)
    make_title_bar(slide, prs, "MLflow — Rastreamento de Experimentos",
                   "Reprodutibilidade · Comparação de Políticas · Registro de Modelos")

    # Section label
    lbl = slide.shapes.add_textbox(Inches(0.45), Inches(1.15), W - Inches(0.9), Inches(0.3))
    lp = lbl.text_frame.paragraphs[0]
    lp.text = "OBSERVABILIDADE & GOVERNANÇA DE ML"
    lp.font.size = Pt(9)
    lp.font.bold = True
    lp.font.color.rgb = GOLD
    lp.font.name = "Inter"

    col_w = (W - Inches(1.2)) / 3
    col_gap = Inches(0.15)
    tops = Inches(1.5)
    card_h = Inches(2.5)

    add_bullet_card(slide,
        Inches(0.45), tops, col_w, card_h,
        "Experiment Tracking", [
            "Cada run = 1 política (Thompson, LinUCB, Nilos-UCB, Baseline)",
            "Parâmetros: horizon, context_dim, alpha, beta",
            "Métricas: cumulative_reward, regret_ratio, conversion_rate",
            "Artefatos: modelo serializado + feature importance",
        ], accent=CYAN)

    add_bullet_card(slide,
        Inches(0.45) + col_w + col_gap, tops, col_w, card_h,
        "Model Registry", [
            "Modelo campeão promovido para 'Production'",
            "Versões anteriores em 'Staging' para rollback",
            "A/B testing via shadow traffic entre versões",
            "API carrega sempre a versão 'Production' ativa",
        ], accent=GREEN)

    add_bullet_card(slide,
        Inches(0.45) + 2*(col_w + col_gap), tops, col_w, card_h,
        "Métricas Comparadas", [
            "Reward acumulado em 6.000 rounds de simulação",
            "Regret ratio: LinUCB 5.1% vs Baseline 100%",
            "Lift vs baseline: LinUCB +66.6% de reward",
            "Convergência observada a partir de round ~800",
        ], accent=GOLD)

    # Metrics table visual
    tbl_top = Inches(4.1)
    tbl_left = Inches(0.45)
    tbl_w = W - Inches(0.9)

    headers = ["Política", "Reward Acum.", "Regret Ratio", "Conversão", "Lift vs Baseline"]
    rows_data = [
        ["LinUCB ★",     "R$ 424.820", "5,1%",   "9,1%",  "+66,6%"],
        ["Thompson",     "R$ 351.200", "17,4%",  "7,8%",  "+37,9%"],
        ["Nilos-UCB",    "R$ 330.100", "22,5%",  "7,3%",  "+29,9%"],
        ["Baseline",     "R$ 254.990", "100,0%", "5,5%",  "—"],
    ]
    row_h = Inches(0.38)
    col_widths = [tbl_w * f for f in [0.22, 0.19, 0.17, 0.17, 0.25]]

    # header row
    x = tbl_left
    for i, hdr in enumerate(headers):
        hbox = slide.shapes.add_shape(1, x, tbl_top, col_widths[i], row_h)
        hbox.fill.solid()
        hbox.fill.fore_color.rgb = VIOLET
        hbox.line.fill.background()
        txb = slide.shapes.add_textbox(x + Pt(4), tbl_top + Pt(4), col_widths[i] - Pt(8), row_h - Pt(8))
        p = txb.text_frame.paragraphs[0]
        p.text = hdr
        p.font.bold = True
        p.font.size = Pt(10)
        p.font.color.rgb = WHITE
        p.font.name = "Inter"
        p.alignment = PP_ALIGN.CENTER
        x += col_widths[i]

    # data rows
    row_colors = [
        RGBColor(0x03, 0x18, 0x3A),  # LinUCB – highlight
        RGBColor(0x03, 0x12, 0x2C),
        RGBColor(0x03, 0x12, 0x2C),
        RGBColor(0x02, 0x0E, 0x22),
    ]
    accent_cols = {0: [GOLD, MUTED, MUTED, MUTED]}  # first row gold text for champion

    for ri, row in enumerate(rows_data):
        x = tbl_left
        for ci, val in enumerate(row):
            rbox = slide.shapes.add_shape(1, x, tbl_top + row_h*(ri+1), col_widths[ci], row_h)
            rbox.fill.solid()
            rbox.fill.fore_color.rgb = row_colors[ri]
            rbox.line.color.rgb = RGBColor(0x0D, 0x1F, 0x42)
            rbox.line.width = Pt(0.5)
            txb = slide.shapes.add_textbox(x + Pt(4), tbl_top + row_h*(ri+1) + Pt(4),
                                           col_widths[ci] - Pt(8), row_h - Pt(8))
            p = txb.text_frame.paragraphs[0]
            p.text = val
            p.font.size = Pt(10)
            p.font.name = "Inter"
            p.alignment = PP_ALIGN.CENTER
            if ri == 0:
                p.font.color.rgb = GOLD
                p.font.bold = True
            elif ci == 0:
                p.font.color.rgb = MUTED
            else:
                p.font.color.rgb = TEXT
            x += col_widths[ci]

    # footer note
    note = slide.shapes.add_textbox(Inches(0.45), H - Inches(0.45), W - Inches(0.9), Inches(0.3))
    np = note.text_frame.paragraphs[0]
    np.text = "★ LinUCB elegido como política campeã via MLflow Model Registry  ·  Simulação em 6.000 rounds com dados reais Kaggle Bank Marketing"
    np.font.size = Pt(8.5)
    np.font.color.rgb = MUTED
    np.font.name = "Inter"
    np.alignment = PP_ALIGN.CENTER


def add_api_slide(prs: Presentation):
    """Add a slide documenting the FastAPI decision service."""
    slide_layout = prs.slide_layouts[0]  # only layout available
    slide = prs.slides.add_slide(slide_layout)
    W, H = prs.slide_width, prs.slide_height

    update_slide_background(slide, PANEL)
    make_title_bar(slide, prs, "API REST — Serviço de Decisão",
                   "FastAPI · 7 endpoints · auditável · localhost:8000/docs")

    lbl = slide.shapes.add_textbox(Inches(0.45), Inches(1.15), W - Inches(0.9), Inches(0.3))
    lp = lbl.text_frame.paragraphs[0]
    lp.text = "CONTRATO DE DECISÃO EM TEMPO REAL"
    lp.font.size = Pt(9)
    lp.font.bold = True
    lp.font.color.rgb = GOLD
    lp.font.name = "Inter"

    endpoints = [
        ("GET",  "/health",               GREEN,  "Liveness + readiness · status, policy_loaded, feature_store"),
        ("GET",  "/policy",               CYAN,   "Política ativa · nome, versão, métricas de treino"),
        ("GET",  "/offers",               CYAN,   "Catálogo de 6 ofertas · id, margem, suitability_tier"),
        ("POST", "/decide",               GOLD,   "Contexto → decisão auditável · arm_id, estimates, reason_codes"),
        ("POST", "/assistant/explain",    AMBER,  "Decide + LLM+RAG · resposta em linguagem natural + citações"),
        ("GET",  "/metrics",              VIOLET, "Matriz de comparação de políticas · reward, regret, lift"),
        ("GET",  "/metrics/regret-curve", VIOLET, "Curvas de regret acumulado para visualização"),
    ]

    row_h = Inches(0.52)
    ep_top = Inches(1.55)
    label_w = Inches(0.72)
    path_w = Inches(2.3)
    desc_w = W - Inches(0.9) - label_w - path_w - Inches(0.1)

    for i, (method, path, color, desc) in enumerate(endpoints):
        top = ep_top + row_h * i
        # method badge
        badge = slide.shapes.add_shape(1, Inches(0.45), top + Inches(0.07), label_w, row_h - Inches(0.14))
        badge.fill.solid()
        badge.fill.fore_color.rgb = color
        badge.line.fill.background()
        mtxb = slide.shapes.add_textbox(Inches(0.45), top + Inches(0.07), label_w, row_h - Inches(0.14))
        mp = mtxb.text_frame.paragraphs[0]
        mp.text = method
        mp.font.bold = True
        mp.font.size = Pt(10)
        mp.font.color.rgb = RGBColor(0x00, 0x00, 0x00) if color in (GOLD, AMBER, GREEN) else WHITE
        mp.font.name = "Inter"
        mp.alignment = PP_ALIGN.CENTER

        # path
        path_box = slide.shapes.add_textbox(Inches(0.45) + label_w + Inches(0.08), top,
                                            path_w, row_h)
        pp = path_box.text_frame.paragraphs[0]
        pp.text = path
        pp.font.bold = True
        pp.font.size = Pt(11)
        pp.font.color.rgb = TEXT
        pp.font.name = "Courier New"
        pp.alignment = PP_ALIGN.LEFT

        # description
        desc_box = slide.shapes.add_textbox(Inches(0.45) + label_w + path_w + Inches(0.18), top,
                                            desc_w, row_h)
        dp = desc_box.text_frame.paragraphs[0]
        dp.text = desc
        dp.font.size = Pt(10)
        dp.font.color.rgb = MUTED
        dp.font.name = "Inter"

        # separator line
        if i < len(endpoints) - 1:
            line_top = top + row_h - Inches(0.04)
            sep = slide.shapes.add_shape(1, Inches(0.45), line_top, W - Inches(0.9), Pt(0.5))
            sep.fill.solid()
            sep.fill.fore_color.rgb = RGBColor(0x0D, 0x1F, 0x42)
            sep.line.fill.background()

    # Decision flow diagram (simplified text representation)
    flow_top = ep_top + row_h * len(endpoints) + Inches(0.15)
    flow_box = slide.shapes.add_shape(1, Inches(0.45), flow_top, W - Inches(0.9), Inches(0.9))
    flow_box.fill.solid()
    flow_box.fill.fore_color.rgb = RGBColor(0x02, 0x0E, 0x22)
    flow_box.line.color.rgb = RGBColor(0x0D, 0x1F, 0x42)
    flow_box.line.width = Pt(0.8)

    ftxb = slide.shapes.add_textbox(Inches(0.65), flow_top + Inches(0.1), W - Inches(1.3), Inches(0.7))
    ftf = ftxb.text_frame
    ftf.word_wrap = False

    lines = [
        ("Fluxo de Decisão: ", GOLD, True),
        ("POST /decide → Feature Extraction → Eligibility Guard → Bandit Policy (Thompson/LinUCB/Nilos-UCB) → Reward Estimate → Audit Log → DecisionOut", MUTED, False),
    ]
    fp = ftf.paragraphs[0]
    for txt, col, bold in lines:
        run = fp.add_run()
        run.text = txt
        run.font.size = Pt(10)
        run.font.color.rgb = col
        run.font.bold = bold
        run.font.name = "Inter"

    # footer
    note = slide.shapes.add_textbox(Inches(0.45), H - Inches(0.45), W - Inches(0.9), Inches(0.3))
    np = note.text_frame.paragraphs[0]
    np.text = "Swagger UI automático em http://localhost:8000/docs  ·  Decisões auditáveis em audit_log.jsonl  ·  FIAP 7MLET Grupo 74"
    np.font.size = Pt(8.5)
    np.font.color.rgb = MUTED
    np.font.name = "Inter"
    np.alignment = PP_ALIGN.CENTER


def recolor_presentation(prs: Presentation):
    """Walk all slides and remap old colors to new palette."""
    for slide in prs.slides:
        # background
        try:
            bg = slide.background.fill
            if bg.type is not None:
                try:
                    old_c = bg.fore_color.rgb
                    new_c = remap_color(old_c)
                    if new_c:
                        bg.fore_color.rgb = new_c
                except Exception:
                    pass
        except Exception:
            pass

        for shape in slide.shapes:
            process_shape(shape)


def main():
    src = Path(__file__).resolve().parents[1] / "docs" / "Adaptive-Offers-Pitch-Grupo74.pptx"
    dst = src  # overwrite in place

    if not src.exists():
        print(f"ERROR: PPTX not found at {src}", file=sys.stderr)
        sys.exit(1)

    print(f"Loading {src} …")
    prs = Presentation(str(src))

    print("Recoloring existing slides …")
    recolor_presentation(prs)

    print("Adding MLflow slide …")
    add_mlflow_slide(prs)

    print("Adding API documentation slide …")
    add_api_slide(prs)

    prs.save(str(dst))
    print(f"✓ Saved to {dst}")
    print(f"  Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
