"""Insere o logo da Adaptive Offers (desenhado com shapes nativas do PowerPoint)
no slide de título do pitch, de forma IDEMPOTENTE.

python-pptx não renderiza SVG, mas o logo é puramente geométrico — então o
recriamos com autoshapes vetoriais (badge + 3 barras + linha + ponto dourado),
preservando a marca e mantendo tudo editável no PowerPoint.

Reaplicar é seguro: se o logo já existir no slide (shape nomeado ``ao_logo_badge``),
o script pula. Não adiciona/duplica nenhum slide.

Uso:  python scripts/add_logo_to_pptx.py
"""

from __future__ import annotations

import shutil
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.util import Emu, Inches, Pt

# ── Paleta da marca ──────────────────────────────────────────────────────────
PANEL   = RGBColor(0x03, 0x0D, 0x24)
VIOLET  = RGBColor(0x00, 0x33, 0xCC)
CYAN    = RGBColor(0x1A, 0x6F, 0xFF)
DIMBLUE = RGBColor(0x14, 0x3A, 0x78)   # barra mais baixa (mimetiza opacidade)
GOLD    = RGBColor(0xFF, 0xC2, 0x00)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)

LOGO_TAG = "ao_logo_badge"


def add_logo(slide, left: int, top: int, size: int) -> None:
    """Desenha o logo num quadrado (left, top, size) em EMU, mapeando o viewBox 48×48."""
    sc = size / 48.0
    def X(v: float) -> int: return int(left + v * sc)
    def Y(v: float) -> int: return int(top + v * sc)
    def D(v: float) -> int: return int(v * sc)

    # badge (fundo navy, borda ciano)
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, X(2), Y(2), D(44), D(44))
    badge.name = LOGO_TAG
    badge.adjustments[0] = 0.28
    badge.fill.solid(); badge.fill.fore_color.rgb = PANEL
    badge.line.color.rgb = CYAN; badge.line.width = Pt(1.5)
    badge.shadow.inherit = False

    # 3 barras ascendentes (braços do bandit)
    for bx, by, bw, bh, col in [
        (11.5, 28, 5.2, 9,  DIMBLUE),
        (20.4, 22, 5.2, 15, CYAN),
        (29.3, 15, 5.2, 22, VIOLET),
    ]:
        r = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, X(bx), Y(by), D(bw), D(bh))
        r.adjustments[0] = 0.35
        r.fill.solid(); r.fill.fore_color.rgb = col
        r.line.fill.background(); r.shadow.inherit = False

    # linha de aprendizado (2 segmentos brancos)
    for x1, y1, x2, y2 in [(11, 30.5, 23, 24), (23, 24, 32, 12.5)]:
        ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, X(x1), Y(y1), X(x2), Y(y2))
        ln.line.color.rgb = WHITE; ln.line.width = Pt(1.6)
        ln.shadow.inherit = False

    # ponto dourado (braço escolhido)
    d = D(7.4)
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, X(32 - 3.7), Y(12.5 - 3.7), d, d)
    dot.fill.solid(); dot.fill.fore_color.rgb = GOLD
    dot.line.color.rgb = WHITE; dot.line.width = Pt(1.1)
    dot.shadow.inherit = False


def main() -> None:
    root = Path(__file__).resolve().parents[1]
    src = root / "docs" / "Adaptive-Offers-Pitch-Grupo64.pptx"
    if not src.exists():
        print(f"ERRO: PPTX não encontrado em {src}", file=sys.stderr)
        sys.exit(1)

    prs = Presentation(str(src))
    slide = prs.slides[0]

    if any(getattr(sh, "name", "") == LOGO_TAG for sh in slide.shapes):
        print("Logo já presente no slide de título — nada a fazer (idempotente).")
        return

    # canto superior direito, fora da área de texto do título
    size = Inches(0.82)
    left = prs.slide_width - size - Inches(0.42)
    top = Inches(0.34)
    add_logo(slide, int(left), int(top), int(size))

    prs.save(str(src))
    print(f"✓ Logo inserido em {src}")

    # espelha no artefato de build, se existir
    build = root / ".pptx-build" / "Adaptive-Offers-Pitch-Grupo64.pptx"
    if build.exists():
        shutil.copy2(src, build)
        print(f"✓ Copiado para {build}")

    print(f"  Total de slides (inalterado): {len(prs.slides)}")


if __name__ == "__main__":
    main()
