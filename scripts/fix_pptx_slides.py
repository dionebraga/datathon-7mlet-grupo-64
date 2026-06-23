"""
Remove slides 20-21 (ruim) e recria slides limpos para MLflow e API.
Versão 2: limpa placeholders herdados do layout e adiciona retângulo
de fundo escuro explícito como primeira shape.
"""
from __future__ import annotations
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

PANEL  = RGBColor(0x03, 0x0D, 0x24)
PANEL2 = RGBColor(0x01, 0x0A, 0x1A)
GRID_C = RGBColor(0x0D, 0x1F, 0x42)
VIOLET = RGBColor(0x00, 0x33, 0xCC)
CYAN   = RGBColor(0x1A, 0x6F, 0xFF)
GREEN  = RGBColor(0x1A, 0x9E, 0x1A)
GOLD   = RGBColor(0xFF, 0xC2, 0x00)
AMBER  = RGBColor(0xFF, 0x9A, 0x00)
RED    = RGBColor(0xE8, 0x40, 0x00)
TEXT   = RGBColor(0xED, 0xED, 0xED)
MUTED  = RGBColor(0x88, 0x99, 0xBB)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
BLACK  = RGBColor(0x00, 0x00, 0x00)


def delete_slide(prs: Presentation, index: int) -> None:
    sldId_lst  = prs.slides._sldIdLst
    sldId_elem = sldId_lst[index]
    r_id       = sldId_elem.get(qn("r:id"))
    prs.part.drop_rel(r_id)
    sldId_lst.remove(sldId_elem)


def clear_placeholders(slide) -> None:
    for ph in list(slide.placeholders):
        sp = ph._element
        sp.getparent().remove(sp)


def rect(slide, left, top, w, h, fill: RGBColor,
         border: RGBColor | None = None, bw: float = 1.0):
    s = slide.shapes.add_shape(1, left, top, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border
        s.line.width = Pt(bw)
    else:
        s.line.fill.background()
    return s


def txt(slide, left, top, w, h, text: str, size: float = 11,
        color: RGBColor = TEXT, bold: bool = False,
        align: PP_ALIGN = PP_ALIGN.LEFT, wrap: bool = False,
        italic: bool = False, name: str = "Calibri"):
    txb = slide.shapes.add_textbox(left, top, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.text = text
    p.font.size   = Pt(size)
    p.font.color.rgb = color
    p.font.bold   = bold
    p.font.italic = italic
    p.font.name   = name
    p.alignment   = align


def header(slide, prs, title: str, subtitle: str = "") -> None:
    W = prs.slide_width
    # full-width gold bar
    rect(slide, 0, 0, W, Inches(0.06), GOLD)
    # dark strip
    rect(slide, 0, Inches(0.06), W, Inches(0.74), PANEL2)
    # title
    txt(slide, Inches(0.45), Inches(0.08), W - Inches(0.9), Inches(0.46),
        title, size=24, bold=True, color=TEXT)
    if subtitle:
        txt(slide, Inches(0.45), Inches(0.54), W - Inches(0.9), Inches(0.24),
            subtitle, size=11, color=MUTED)


def section_lbl(slide, left, top, w, text: str) -> None:
    txt(slide, left, top, w, Inches(0.25), text, size=8, bold=True, color=GOLD)


# ─── MLflow slide ────────────────────────────────────────────────────────────
def mlflow_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    W, H  = prs.slide_width, prs.slide_height
    clear_placeholders(slide)

    # ① Full-slide background rectangle (must be FIRST shape)
    rect(slide, 0, 0, W, H, PANEL)

    # Header
    header(slide, prs,
           "MLflow — Rastreamento de Experimentos",
           "Reprodutibilidade · Comparação de Políticas · Registro de Modelos")

    section_lbl(slide, Inches(0.45), Inches(0.92), W - Inches(0.9),
                "OBSERVABILIDADE E GOVERNANÇA DE ML")

    # ── 3 cards ──────────────────────────────────────────────────────────────
    GAP = Inches(0.20)
    CW  = (W - Inches(0.9) - 2 * GAP) / 3
    CT  = Inches(1.18)
    CH  = Inches(2.40)

    cards = [
        ("Experiment Tracking", CYAN, [
            "1 run = 1 politica simulada",
            "params: horizon, alpha, beta",
            "metricas: reward, regret",
            "artefatos: modelo + features",
        ]),
        ("Model Registry", GREEN, [
            "LinUCB promovido a Production",
            "versoes anteriores em Staging",
            "rollback disponivel",
            "API usa versao Production",
        ]),
        ("Resultados-chave", GOLD, [
            "LinUCB: regret ratio 5,1%",
            "Lift +66,6% vs baseline",
            "Convergencia: round ~800",
            "6.000 rounds, dados reais",
        ]),
    ]

    for i, (ctitle, accent, bullets) in enumerate(cards):
        lft = Inches(0.45) + i * (CW + GAP)
        rect(slide, lft, CT, CW, CH, PANEL2, border=accent, bw=1.5)
        # card title bar
        rect(slide, lft, CT, CW, Inches(0.38), accent)
        txt(slide, lft + Inches(0.12), CT + Inches(0.05),
            CW - Inches(0.24), Inches(0.30),
            ctitle, size=12, bold=True,
            color=BLACK if accent in (GOLD,) else WHITE)
        # bullets
        for j, b in enumerate(bullets):
            txt(slide,
                lft + Inches(0.16),
                CT + Inches(0.45) + j * Inches(0.46),
                CW - Inches(0.32),
                Inches(0.42),
                f"  {b}", size=10, color=TEXT)

    # ── Table ────────────────────────────────────────────────────────────────
    TT = CT + CH + Inches(0.28)
    TW = W - Inches(0.9)
    RH = Inches(0.40)
    TL = Inches(0.45)

    headers = ["Politica", "Reward Acum.", "Regret", "Conversao", "Lift vs Baseline"]
    fracs   = [0.22, 0.20, 0.14, 0.14, 0.30]
    cols    = [TW * f for f in fracs]

    x = TL
    for hi, h in enumerate(headers):
        rect(slide, x, TT, cols[hi], RH, VIOLET)
        txt(slide, x + Inches(0.05), TT + Inches(0.07),
            cols[hi] - Inches(0.10), RH - Inches(0.10),
            h, size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        x += cols[hi]

    rows_data = [
        ("LinUCB (campeao)", "R$ 424.820", "5,1%", "9,1%", "+66,6%", True),
        ("Thompson",         "R$ 351.200", "17,4%","7,8%", "+37,9%", False),
        ("Nilos-UCB",        "R$ 330.100", "22,5%","7,3%", "+29,9%", False),
        ("Baseline (greedy)","R$ 254.990", "100%", "5,5%", "---",    False),
    ]
    row_bgs = [RGBColor(0x03,0x18,0x3A), RGBColor(0x03,0x10,0x28),
               RGBColor(0x03,0x10,0x28), RGBColor(0x02,0x0C,0x1E)]

    for ri, (pol, rwd, rgt, cnv, lft, champ) in enumerate(rows_data):
        rrow = TT + RH * (ri + 1)
        x    = TL
        for ci, cell in enumerate([pol, rwd, rgt, cnv, lft]):
            rect(slide, x, rrow, cols[ci], RH, row_bgs[ri], border=GRID_C, bw=0.4)
            tc = GOLD if champ else (MUTED if ci == 0 else TEXT)
            txt(slide, x + Inches(0.05), rrow + Inches(0.08),
                cols[ci] - Inches(0.10), RH - Inches(0.10),
                cell, size=9, bold=champ, color=tc, align=PP_ALIGN.CENTER)
            x += cols[ci]

    # footer
    txt(slide, TL, H - Inches(0.36), W - Inches(0.9), Inches(0.28),
        "LinUCB eleito campeao via MLflow Model Registry  |  "
        "Simulacao em 6.000 rounds  |  Bank Marketing Dataset  |  41.188 clientes reais",
        size=8, color=MUTED, align=PP_ALIGN.CENTER)


# ─── API slide ────────────────────────────────────────────────────────────────
def api_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    W, H  = prs.slide_width, prs.slide_height
    clear_placeholders(slide)

    # ① Full-slide background
    rect(slide, 0, 0, W, H, PANEL)

    header(slide, prs,
           "API REST — Servico de Decisao",
           "FastAPI · 7 endpoints · auditavel · localhost:8000/docs")

    section_lbl(slide, Inches(0.45), Inches(0.92), W - Inches(0.9),
                "CONTRATO DE DECISAO EM TEMPO REAL")

    endpoints = [
        ("GET",  "/health",               GREEN,  "Liveness + readiness"),
        ("GET",  "/policy",               CYAN,   "Politica ativa - nome, versao, metricas"),
        ("GET",  "/offers",               CYAN,   "Catalogo de 6 ofertas - margem, suitability"),
        ("POST", "/decide",               GOLD,   "Contexto -> decisao auditavel + reason_codes"),
        ("POST", "/assistant/explain",    AMBER,  "Decide + LLM+RAG - resposta em linguagem natural"),
        ("GET",  "/metrics",              VIOLET, "Matriz de comparacao de politicas"),
        ("GET",  "/metrics/regret-curve", VIOLET, "Curvas de regret acumulado para visualizacao"),
    ]

    RH      = Inches(0.50)
    ET      = Inches(1.18)
    BW      = Inches(0.68)
    PW      = Inches(2.40)
    DW      = W - Inches(0.9) - BW - PW - Inches(0.20)

    for i, (method, path, color, desc) in enumerate(endpoints):
        top = ET + RH * i

        # separator
        if i > 0:
            rect(slide, Inches(0.45), top, W - Inches(0.9), Pt(0.5), GRID_C)

        mid = top + Inches(0.06)
        bh  = RH - Inches(0.12)

        # method badge
        rect(slide, Inches(0.45), mid, BW, bh, color)
        bc = BLACK if color in (GOLD, AMBER, GREEN) else WHITE
        txt(slide, Inches(0.45), mid, BW, bh,
            method, size=10, bold=True, color=bc, align=PP_ALIGN.CENTER)

        # path
        txt(slide, Inches(0.45) + BW + Inches(0.12), top + Inches(0.08),
            PW, RH - Inches(0.10),
            path, size=11, bold=True, color=TEXT, name="Consolas")

        # description
        txt(slide, Inches(0.45) + BW + PW + Inches(0.18), top + Inches(0.10),
            DW, RH - Inches(0.10),
            desc, size=10, color=MUTED)

    # Decision flow box
    FT = ET + RH * len(endpoints) + Inches(0.22)
    FH = Inches(0.80)
    rect(slide, Inches(0.45), FT, W - Inches(0.9), FH, PANEL2, border=GRID_C)

    txt(slide, Inches(0.65), FT + Inches(0.10),
        W - Inches(1.3), Inches(0.24),
        "FLUXO DE DECISAO", size=8, bold=True, color=GOLD)

    txt(slide, Inches(0.65), FT + Inches(0.36),
        W - Inches(1.3), Inches(0.32),
        "POST /decide  ->  Feature Extraction  ->  Eligibility Guard  ->  "
        "Bandit Policy  ->  Reward Estimate  ->  Audit Log  ->  DecisionOut",
        size=10, color=MUTED)

    # footer
    txt(slide, Inches(0.45), H - Inches(0.36), W - Inches(0.9), Inches(0.28),
        "Swagger UI em http://localhost:8000/docs  |  "
        "Decisoes auditaveis em audit_log.jsonl  |  FIAP 7MLET Grupo 74",
        size=8, color=MUTED, align=PP_ALIGN.CENTER)


# ─── Main ────────────────────────────────────────────────────────────────────
def main() -> None:
    src = Path(__file__).resolve().parents[1] / "docs" / "Adaptive-Offers-Pitch-Grupo74.pptx"
    if not src.exists():
        sys.exit(f"ERROR: {src}")

    prs = Presentation(str(src))
    n   = len(prs.slides)
    print(f"Loaded: {n} slides")

    # Remove last 2 bad slides
    while len(prs.slides) > 19:
        delete_slide(prs, len(prs.slides) - 1)
        print(f"  Removed slide, now {len(prs.slides)}")

    print("Adding MLflow slide...")
    mlflow_slide(prs)

    print("Adding API slide...")
    api_slide(prs)

    prs.save(str(src))
    build = src.parent.parent / ".pptx-build" / src.name
    build.parent.mkdir(exist_ok=True)
    prs.save(str(build))
    print(f"Done: {len(prs.slides)} slides -> {src}")


if __name__ == "__main__":
    main()
